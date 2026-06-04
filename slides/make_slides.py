#!/usr/bin/env python3
"""Build informal graph slides for the cvc5 cycle/minimal work.

Graphs are drawn with native PowerPoint shapes so they remain editable in
Google Slides. Run:  python3 make_slides.py   ->  cycle_slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- palette: one color per relation -------------------------------------
RED = RGBColor(0xD3, 0x2F, 0x2F)   # R1
BLUE = RGBColor(0x15, 0x65, 0xC0)  # R2
GREEN = RGBColor(0x2E, 0x7D, 0x32) # R3
BLACK = RGBColor(0x21, 0x21, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)  # dimmed / not-explored

NODE_DIA = Inches(0.75)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen (Google Slides default)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]       # fully blank layout, no placeholders


def add_node(slide, label, cx, cy):
    """Draw a circular node centered at (cx, cy) inches; return center point."""
    left = Inches(cx) - NODE_DIA // 2
    top = Inches(cy) - NODE_DIA // 2
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, NODE_DIA, NODE_DIA)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = BLACK
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = BLACK
    return (Inches(cx), Inches(cy))


def _set_arrowhead(connector):
    """Add a triangular arrowhead at the connector's end."""
    ln = connector.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'),
                          {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)


def add_edge(slide, p_from, p_to, color, offset=0.0):
    """Draw a colored directed arrow from p_from to p_to (center points).

    Endpoints are pulled back to the node rim; `offset` shifts the line
    perpendicularly (inches) so anti-parallel edges don't overlap.
    """
    import math
    x1, y1 = p_from[0], p_from[1]
    x2, y2 = p_to[0], p_to[1]
    dx, dy = (x2 - x1), (y2 - y1)
    dist = math.hypot(dx, dy) or 1
    ux, uy = dx / dist, dy / dist          # unit direction
    px, py = -uy, ux                       # unit perpendicular
    r = NODE_DIA / 2 + Emu(20000)          # pull back to rim + small gap
    off = Inches(offset)
    sx = x1 + Emu(int(ux * r)) + Emu(int(px * off))
    sy = y1 + Emu(int(uy * r)) + Emu(int(py * off))
    ex = x2 - Emu(int(ux * r)) + Emu(int(px * off))
    ey = y2 - Emu(int(uy * r)) + Emu(int(py * off))
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)
    conn.line.color.rgb = color
    conn.line.width = Pt(2.5)
    _set_arrowhead(conn)
    return conn


def add_relation_text(slide):
    """Top-left text block: the three relations + the acyclicity assertion."""
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3),
                                   Inches(5.5), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        ("R1 = {(a, b)}", RED),
        ("R2 = {(b, a)}", BLUE),
        ("R3 = {(c, d)}", GREEN),
    ]
    for i, (txt, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.name = 'Consolas'
        r.font.color.rgb = col
    # assertion underneath
    p = tf.add_paragraph()
    p.space_before = Pt(18)
    r = p.add_run()
    r.text = "¬ acyclic(R1, R2, R3)"   # ¬ acyclic(...)
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = BLACK

    # clarifying comment underneath the assertion
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    r = p.add_run()
    r.text = "i.e., R1 ∪ R2 ∪ R3 contains a cycle"
    r.font.size = Pt(18)
    r.font.italic = True
    r.font.color.rgb = GRAY


# ---- Rules reference slide (embedded images) ------------------------------
# Add the other two rule images here as they arrive; they stack vertically.
RULE_BLUE = RGBColor(0x12, 0x12, 0xCC)

# rule 1 (InstCycle) is now typeset as editable text; 2 and 3 stay as images
IMAGE_RULES = [
    "rules/rule2_splitcyclelen.png",
    "rules/rule3_unrollcycle.png",
]


def _text_line(slide, text, left, top, width, height, size,
               color=RULE_BLUE, bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def draw_rule(slide, label, premises, conclusions, top, bar_w, size=18):
    """Typeset a general inference rule (premises / bar / conclusions).

    Returns the y-coordinate of the rule's bottom.
    """
    body_x = Inches(2.5)
    body_w = prs.slide_width - Inches(0.7) - body_x
    line_h = Inches(0.36)

    # premises (above the bar)
    y = top
    for line in premises:
        _text_line(slide, line, body_x, y, body_w, line_h, size)
        y += line_h

    # inference bar
    bar_y = y + Inches(0.05)
    bar_left = body_x + (body_w - bar_w) // 2
    bar = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     bar_left, bar_y, bar_left + bar_w, bar_y)
    bar.line.color.rgb = RULE_BLUE
    bar.line.width = Pt(1.5)

    # conclusions (below the bar)
    y = bar_y + Inches(0.08)
    for line in conclusions:
        _text_line(slide, line, body_x, y, body_w, line_h, size)
        y += line_h

    # label, vertically centered on the rule
    label_y = (top + y) // 2 - Inches(0.18)
    _text_line(slide, label, Inches(0.55), label_y, Inches(2.0), line_h, size,
               align=PP_ALIGN.LEFT)
    return y


def add_rules_slide():
    slide = prs.slides.add_slide(BLANK)

    gap = Inches(0.42)
    # Rule 1 — InstCycle
    bottom = draw_rule(
        slide, "InstCycle",
        ["¬acyclic(R₁,…,Rₖ) ∈ Δ*        ((R₁,…,Rₖ), _, _) ∉ C"],
        ["C := C ∪ {((R₁,…,Rₖ), s, 1)}",
         "Δ := Δ ∪ {len(s) > 1, minimal((R₁,…,Rₖ), s)}"],
        Inches(0.55), Inches(7.8))

    # Rule 2 — SplitCycleLen
    bottom = draw_rule(
        slide, "SplitCycleLen",
        ["((R₁,…,Rₖ), s, cnt) ∈ C"],
        ["Δ := Δ ∪ {cnt < len(s)}      ‖      Δ := Δ ∪ {cnt ≈ len(s)}"],
        bottom + gap, Inches(8.6))

    # Rule 3 — UnrollCycle
    bottom = draw_rule(
        slide, "UnrollCycle",
        ["((R₁,…,Rₖ), s, cnt) ∈ C        cnt < len(s) ∈ Δ*",
         "C′ ≡ C \\ {((R₁,…,Rₖ), s, cnt)} ∪ {((R₁,…,Rₖ), s, cnt+1)}"],
        ["C := C′      Δ := Δ ∪ {(s[cnt-1], s[cnt]) ∈ R₁⁺}",
         "‖   ⋯",
         "‖   C := C′      Δ := Δ ∪ {(s[cnt-1], s[cnt]) ∈ Rₖ⁺}"],
        bottom + gap, Inches(9.6))

    # Rule 4 — ContrMinimal
    draw_rule(
        slide, "ContrMinimal",
        ["((R₁,…,Rₖ), s, cnt) ∈ C        minimal((R₁,…,Rₖ), s)",
         "b ∈ [1, k]        1 ≤ q < r−1 ≤ len(s)        (s[q], s[r]) ∈ R_b⁺"],
        ["unsat"],
        bottom + gap, Inches(9.6))
    return slide


add_rules_slide()   # becomes slide 1 (created first)

# ---- Relations + graph slide ----------------------------------------------
s1 = prs.slides.add_slide(BLANK)

# title
tbox = s1.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.8))
tp = tbox.text_frame.paragraphs[0]
trun = tp.add_run()
trun.text = "Example 1: SAT instance"
trun.font.size = Pt(32)
trun.font.bold = True
trun.font.color.rgb = BLACK

add_relation_text(s1)

# nodes: a<->b 2-cycle on top, c->d below, on the right half of the slide
a = add_node(s1, "a", 8.3, 2.4)
b = add_node(s1, "b", 11.2, 2.4)
c = add_node(s1, "c", 8.3, 5.2)
d = add_node(s1, "d", 11.2, 5.2)

add_edge(s1, a, b, RED, offset=-0.16)    # R1: a -> b (upper line)
add_edge(s1, b, a, BLUE, offset=-0.16)   # R2: b -> a (lower line, mirrored)
add_edge(s1, c, d, GREEN)                # R3: c -> d


# ---- Derivation slide (Example 1) -----------------------------------------
def _para(tf, runs, first=False, space_before=0, space_after=4):
    """Add a paragraph made of (text, color, size, bold) runs."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    for text, color, size, bold in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return p


def _state_box(slide, label, body, x, y, w=Inches(4.9), h=Inches(0.7)):
    """Small rounded box stating a Δ state, e.g. 'Δ₀ = {…}'."""
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid()
    b.fill.fore_color.rgb = WHITE
    b.line.color.rgb = BLACK
    b.line.width = Pt(1.75)
    b.shadow.inherit = False
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    ra = p.add_run()
    ra.text = label
    ra.font.size = Pt(18)
    ra.font.bold = True
    ra.font.color.rgb = BLACK
    rb = p.add_run()
    rb.text = body
    rb.font.size = Pt(18)
    rb.font.color.rgb = BLACK
    return b


def draw_delta_box(slide, label, member_text, color, x, y, w, h,
                   dim=False, hsize=30, bsize=20):
    accent = GRAY if dim else color
    shared = GRAY if dim else BLACK
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = accent
    shp.line.width = Pt(2.0 if dim else 2.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # header
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(10)
    r = p0.add_run()
    r.text = label
    r.font.size = Pt(hsize)
    r.font.bold = True
    r.font.color.rgb = accent
    # shared element
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(5)
    r1 = p1.add_run()
    r1.text = "1 < len(s)"
    r1.font.size = Pt(bsize)
    r1.font.color.rgb = shared
    # branch-specific element
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = member_text
    r2.font.size = Pt(bsize)
    r2.font.color.rgb = accent
    return shp


def add_derivation_slide():
    slide = prs.slides.add_slide(BLANK)
    # title
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1),
                                 Inches(0.8))
    tr = t.text_frame.paragraphs[0].add_run()
    tr.text = "Example 1: derivation"
    tr.font.size = Pt(32)
    tr.font.bold = True
    tr.font.color.rgb = BLACK

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12.0),
                                   Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    HDR, BODY = 25, 21          # font sizes
    blue = RULE_BLUE

    # Step 1 — InstCycle
    _para(tf, [("Step 1 — InstCycle", BLACK, HDR, True)], first=True)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("¬acyclic(R1, R2, R3)", blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("C = {((R1, R2, R3), s, 1)}", blue, BODY, False)])
    _para(tf, [("            ", BLACK, BODY, False),
               ("Δ = {len(s) > 1, minimal(s)}", blue, BODY, False)])

    # Step 2 — SplitCycleLen
    _para(tf, [("Step 2 — SplitCycleLen", BLACK, HDR, True)], space_before=18)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("((R1, R2, R3), s, 1) ∈ C", blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("Δ = Δ ∪ {1 < len(s)}", blue, BODY, False),
               ("      ‖      ", BLACK, BODY, True),
               ("Δ = Δ ∪ {1 = len(s)}", blue, BODY, False)])
    _para(tf, [("            ", BLACK, BODY, False),
               ("(we only explore case 1)", GRAY, BODY - 4, False)])

    # Step 3 — UnrollCycle (from the 1 < len(s) branch)
    _para(tf, [("Step 3 — UnrollCycle", BLACK, HDR, True)], space_before=16)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("((R1, R2, R3), s, 1) ∈ C", blue, BODY, False),
               ("   and   ", BLACK, BODY, False),
               ("1 < len(s)", blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("Δ = Δ ∪ {(s[0], s[1]) ∈ R1⁺}", blue, BODY, False)])
    _para(tf, [("       ‖   ", BLACK, BODY, True),
               ("Δ = Δ ∪ {(s[0], s[1]) ∈ R2⁺}", blue, BODY, False)])
    _para(tf, [("       ‖   ", BLACK, BODY, True),
               ("Δ = Δ ∪ {(s[0], s[1]) ∈ R3⁺}", blue, BODY, False)])

    # state of Δ next to Steps 1 and 2 (right side)
    _state_box(slide, "Δ₀ = ", "{len(s) > 1, minimal(s)}",
               Inches(7.7), Inches(1.32))
    _state_box(slide, "Δ₀′ = ", "{len(s) > 1, minimal(s)}",
               Inches(7.7), Inches(3.18))

    # three resulting branches next to Step 3 (bottom-right, below premise)
    bw, bh, bgap = Inches(2.25), Inches(1.4), Inches(0.15)
    bx0, by = Inches(5.7), Inches(6.0)
    branches = [
        ("Δ₁", "(s[0], s[1]) ∈ R1⁺", RED),
        ("Δ₂", "(s[0], s[1]) ∈ R2⁺", BLUE),
        ("Δ₃", "(s[0], s[1]) ∈ R3⁺", GREEN),
    ]
    for i, (lab, mem, col) in enumerate(branches):
        draw_delta_box(slide, lab, mem, col, bx0 + i * (bw + bgap), by, bw, bh,
                       hsize=17, bsize=12)
    return slide


add_derivation_slide()   # slide 3


# ---- Explore Delta1: Steps 4 & 5 (slide 4) --------------------------------
def add_explore_delta1_slide():
    slide = prs.slides.add_slide(BLANK)
    # title
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1),
                                 Inches(0.65))
    tr = t.text_frame.paragraphs[0].add_run()
    tr.text = "Example 1: exploring Δ₁ — Steps 4 & 5"
    tr.font.size = Pt(28)
    tr.font.bold = True
    tr.font.color.rgb = BLACK

    note = slide.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(8.4),
                                    Inches(0.4))
    nr = note.text_frame.paragraphs[0].add_run()
    nr.text = "Exploring Δ₁ only (Δ₂, Δ₃ omitted)."
    nr.font.size = Pt(15)
    nr.font.italic = True
    nr.font.color.rgb = GRAY

    # left column: Step 4 and Step 5 derivation text
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(8.4),
                                   Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    HDR, BODY = 23, 19
    blue = RULE_BLUE
    # Step 4
    _para(tf, [("Step 4 — SplitCycleLen", BLACK, HDR, True)], first=True)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("((R1, R2, R3), s, 2) ∈ C", blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("2 < len(s)", blue, BODY, False),
               ("   (case 1 — explored)", BLACK, BODY - 4, False)])
    _para(tf, [("       ‖   ", GRAY, BODY, True),
               ("2 = len(s)", GRAY, BODY, False),
               ("   (case 2 — omitted)", GRAY, BODY - 4, False)])
    # Step 5
    _para(tf, [("Step 5 — UnrollCycle", BLACK, HDR, True)], space_before=14)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("((R1, R2, R3), s, 2) ∈ C", blue, BODY, False),
               ("  and  ", BLACK, BODY, False),
               ("2 < len(s)", blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("(s[1], s[2]) ∈ ", BLACK, BODY, False),
               ("R1⁺", RED, BODY, False),
               ("  ‖  ", BLACK, BODY, True),
               ("R2⁺", BLUE, BODY, False),
               ("  ‖  ", BLACK, BODY, True),
               ("R3⁺", GREEN, BODY, False)])

    # Δ₁′ (result of Step 4) as a box, on the right next to Step 4
    draw_delta_box_multi(slide, "Δ₁′",
                         [("2 < len(s)", BLACK), ("(s[0], s[1]) ∈ R1⁺", RED)],
                         RED, Inches(8.8), Inches(1.7), Inches(4.0),
                         Inches(1.45), hsize=20, bsize=15)

    # Step-5 result branches in a horizontal row at the bottom (Δ₁₂ active)
    bw, bh, bgap = Inches(3.9), Inches(1.95), Inches(0.28)
    bx0, by = Inches(0.55), Inches(4.95)
    common = [("2 < len(s)", BLACK), ("(s[0], s[1]) ∈ R1⁺", RED)]
    specs = [
        ("Δ₁₁", ("(s[1], s[2]) ∈ R1⁺", RED), RED, True),
        ("Δ₁₂", ("(s[1], s[2]) ∈ R2⁺", BLUE), BLUE, False),
        ("Δ₁₃", ("(s[1], s[2]) ∈ R3⁺", GREEN), GREEN, True),
    ]
    for i, (lab, extra, col, dim) in enumerate(specs):
        draw_delta_box_multi(slide, lab, common + [extra], col,
                             bx0 + i * (bw + bgap), by, bw, bh,
                             dim=dim, hsize=18, bsize=13)

    cap = slide.shapes.add_textbox(Inches(0.6), by + bh + Inches(0.05),
                                   Inches(12.1), Inches(0.4))
    c = cap.text_frame.paragraphs[0]
    c.alignment = PP_ALIGN.CENTER
    cr = c.add_run()
    cr.text = "We explore Δ₁₂ only (Δ₁₁, Δ₁₃ omitted)."
    cr.font.size = Pt(14)
    cr.font.italic = True
    cr.font.color.rgb = GRAY
    return slide


# ---- multi-line Δ box helper (used by the merged slide 4) -----------------
def draw_delta_box_multi(slide, label, lines, accent, x, y, w, h,
                         dim=False, hsize=22, bsize=16):
    head = GRAY if dim else accent
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = head
    shp.line.width = Pt(2.0 if dim else 2.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(8)
    r = p0.add_run()
    r.text = label
    r.font.size = Pt(hsize)
    r.font.bold = True
    r.font.color.rgb = head
    for txt, col in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(3)
        rr = p.add_run()
        rr.text = txt
        rr.font.size = Pt(bsize)
        rr.font.color.rgb = GRAY if dim else col
    return shp


add_explore_delta1_slide()   # slide 4 (merged Steps 4 & 5)


# ---- Step 6 + saturation / model (slide 7) --------------------------------
def add_step6_slide():
    slide = prs.slides.add_slide(BLANK)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1),
                                 Inches(0.7))
    tr = t.text_frame.paragraphs[0].add_run()
    tr.text = "Example 1: exploring Δ₁ — Step 6"
    tr.font.size = Pt(30)
    tr.font.bold = True
    tr.font.color.rgb = BLACK

    # reminder: current branch Δ₁₂ (top-right corner)
    draw_delta_box_multi(
        slide, "Δ₁₂",
        [("2 < len(s)", BLACK),
         ("(s[0], s[1]) ∈ R1⁺", RED),
         ("(s[1], s[2]) ∈ R2⁺", BLUE)],
        BLUE, Inches(9.45), Inches(1.05), Inches(3.35), Inches(1.95),
        hsize=18, bsize=13)

    # Step 6 derivation text (this time case 2 is explored)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12.0),
                                   Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    HDR, BODY = 25, 21
    blue = RULE_BLUE
    _para(tf, [("Step 6 — SplitCycleLen", BLACK, HDR, True)], first=True)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("((R1, R2, R3), s, 3) ∈ C", blue, BODY, False)])
    _para(tf, [("       ‖   ", GRAY, BODY, True),
               ("3 < len(s)", GRAY, BODY, False),
               ("    (case 1 — omitted)", GRAY, BODY - 4, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("3 = len(s)", blue, BODY, False),
               ("    (case 2 — explored)", BLACK, BODY - 4, False)])

    # saturation reasoning (bullets)
    bb = slide.shapes.add_textbox(Inches(0.7), Inches(3.15), Inches(12.0),
                                  Inches(2.7))
    btf = bb.text_frame
    btf.word_wrap = True
    SZ = 16
    _para(btf, [("•  All cycle-related rules are now saturated", BLACK, SZ, True),
                ("  — assuming a guard on ", BLACK, SZ, False),
                ("SplitCycleLen", BLACK, SZ, True),
                (" that we haven't reached the end of the cycle (no constraint "
                 "len(s) = n for a constant n in Δ).", BLACK, SZ, False)],
          first=True, space_after=8)
    _para(btf, [("•  We can't apply ", BLACK, SZ, False),
                ("UnrollCycle", BLACK, SZ, True),
                (", whose premise ", BLACK, SZ, False),
                ("cnt < len(s)", blue, SZ, False),
                (" — here ", BLACK, SZ, False),
                ("3 < len(s)", blue, SZ, False),
                (" — is not in Δ, since we have ", BLACK, SZ, False),
                ("3 = len(s)", blue, SZ, False),
                (".", BLACK, SZ, False)], space_after=8)
    _para(btf, [("•  Once we also saturate the original relational-solver "
                 "rules, we can conclude ", BLACK, SZ, False),
                ("sat", GREEN, SZ, True),
                (".", BLACK, SZ, False)], space_after=8)
    _para(btf, [("•  Congruence-closure reasoning concludes ", BLACK, SZ, False),
                ("s[0] = a, s[1] = b, s[2] = a", BLACK, SZ, True),
                (", which gives our model of R.", BLACK, SZ, False)],
          space_after=8)

    # model box
    mb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(2.4), Inches(6.15), Inches(8.5),
                                Inches(0.85))
    mb.fill.solid()
    mb.fill.fore_color.rgb = WHITE
    mb.line.color.rgb = GREEN
    mb.line.width = Pt(2.5)
    mb.shadow.inherit = False
    mtf = mb.text_frame
    mtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    mp = mtf.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    ra = mp.add_run()
    ra.text = "Model:  "
    ra.font.size = Pt(22)
    ra.font.bold = True
    ra.font.color.rgb = GREEN
    rb = mp.add_run()
    rb.text = "s[0] = a,  s[1] = b,  s[2] = a"
    rb.font.size = Pt(22)
    rb.font.color.rgb = BLACK
    return slide


add_step6_slide()   # slide 5


# ===== Example 2: unsat instance ===========================================
def draw_circle(slide, cx, cy, dia, label="", line_color=BLACK,
                fill=WHITE, lw=1.5, fontsize=18):
    left = Inches(cx) - Inches(dia) // 2
    top = Inches(cy) - Inches(dia) // 2
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top,
                                 Inches(dia), Inches(dia))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if label:
        p = shp.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(fontsize)
        r.font.bold = True
        r.font.color.rgb = BLACK
    return (Inches(cx), Inches(cy))


def draw_arrow(slide, p_from, p_to, color, pullback, dotted=False, width=2.0):
    import math
    x1, y1 = p_from
    x2, y2 = p_to
    dx, dy = (x2 - x1), (y2 - y1)
    dist = math.hypot(dx, dy) or 1
    ux, uy = dx / dist, dy / dist
    sx = x1 + Emu(int(ux * pullback))
    sy = y1 + Emu(int(uy * pullback))
    ex = x2 - Emu(int(ux * pullback))
    ey = y2 - Emu(int(uy * pullback))
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, sy, ex, ey)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dotted:
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    _set_arrowhead(conn)
    return conn


def add_ex2_instance_slide():
    import math
    slide = prs.slides.add_slide(BLANK)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1),
                                 Inches(0.7))
    tr = t.text_frame.paragraphs[0].add_run()
    tr.text = "Example 2: unsat instance"
    tr.font.size = Pt(32)
    tr.font.bold = True
    tr.font.color.rgb = BLACK

    # constraints (left)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(6.0),
                                   Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(14)
    r = p.add_run()
    r.text = "∀x : Int.  ¬ ( (x, x) ∈ R⁺ )"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = BLACK
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "¬ acyclic(R)"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = BLACK
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = "i.e., R has a cycle, but R⁺ has no self-loop"
    r.font.size = Pt(16)
    r.font.italic = True
    r.font.color.rgb = GRAY

    # 5-node cycle in R (right), unlabeled nodes
    cx, cy, rad = 10.0, 3.4, 1.45
    nodes = []
    for i in range(5):
        ang = math.radians(90 - i * 72)
        nodes.append(draw_circle(slide, cx + rad * math.cos(ang),
                                 cy - rad * math.sin(ang), 0.5))
    pull = Inches(0.25) + Emu(15000)
    for i in range(5):
        draw_arrow(slide, nodes[i], nodes[(i + 1) % 5], BLACK, pull, width=2.0)
    gcap = slide.shapes.add_textbox(Inches(8.5), Inches(5.05), Inches(3.0),
                                    Inches(0.4))
    gc = gcap.text_frame.paragraphs[0]
    gc.alignment = PP_ALIGN.CENTER
    gr = gc.add_run()
    gr.text = "a cycle in R"
    gr.font.size = Pt(16)
    gr.font.italic = True
    gr.font.color.rgb = GRAY

    # node x with a dotted self-loop labeled R⁺, crossed out (does not exist)
    xc = draw_circle(slide, 2.8, 5.7, 0.55, "x", fontsize=20)
    loop = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.8 - 0.4),
                                  Inches(5.7 - 1.15), Inches(0.8), Inches(0.95))
    loop.fill.background()
    loop.line.color.rgb = RULE_BLUE
    loop.line.width = Pt(2.0)
    lln = loop.line._get_or_add_ln()
    lln.append(lln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    loop.shadow.inherit = False
    # R⁺ label on the loop
    lbl = slide.shapes.add_textbox(Inches(3.35), Inches(4.5), Inches(1.2),
                                   Inches(0.5))
    lr = lbl.text_frame.paragraphs[0].add_run()
    lr.text = "R⁺"
    lr.font.size = Pt(20)
    lr.font.bold = True
    lr.font.color.rgb = RULE_BLUE
    # red X showing the self-loop does not exist
    xmark = slide.shapes.add_textbox(Inches(3.5), Inches(5.2), Inches(1.0),
                                     Inches(0.9))
    xm = xmark.text_frame.paragraphs[0].add_run()
    xm.text = "✗"
    xm.font.size = Pt(44)
    xm.font.bold = True
    xm.font.color.rgb = RED
    xcap = slide.shapes.add_textbox(Inches(1.4), Inches(6.55), Inches(4.5),
                                    Inches(0.5))
    xcp = xcap.text_frame.paragraphs[0]
    xcp.alignment = PP_ALIGN.CENTER
    xr = xcp.add_run()
    xr.text = "no self-loop (x, x) exists in R⁺"
    xr.font.size = Pt(16)
    xr.font.italic = True
    xr.font.color.rgb = GRAY
    return slide


def add_ex2_derivation_slide():
    slide = prs.slides.add_slide(BLANK)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12.1),
                                 Inches(0.7))
    tr = t.text_frame.paragraphs[0].add_run()
    tr.text = "Example 2: derivation"
    tr.font.size = Pt(30)
    tr.font.bold = True
    tr.font.color.rgb = BLACK

    # Steps 1-3 (compact, left)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(7.0),
                                   Inches(3.3))
    tf = box.text_frame
    tf.word_wrap = True
    HDR, BODY = 19, 15
    blue = RULE_BLUE
    _para(tf, [("Step 1 — InstCycle", BLACK, HDR, True)],
          first=True, space_after=2)
    _para(tf, [("   Premise(s):  ", BLACK, BODY, False),
               ("¬acyclic(R)", blue, BODY, False)])
    _para(tf, [("   ⟹  ", BLACK, BODY, False),
               ("C = {((R), s, 1)},   Δ = {len(s) > 1, minimal((R), s)}",
                blue, BODY, False)])
    _para(tf, [("Step 2 — SplitCycleLen", BLACK, HDR, True)],
          space_before=8, space_after=2)
    _para(tf, [("   Premise(s):  ", BLACK, BODY, False),
               ("((R), s, 1) ∈ C", blue, BODY, False)])
    _para(tf, [("   ⟹  ", BLACK, BODY, False),
               ("1 < len(s)", blue, BODY, False),
               ("  ‖  ", BLACK, BODY, True),
               ("1 = len(s)", blue, BODY, False),
               ("   (case 1)", GRAY, BODY - 3, False)])
    _para(tf, [("Step 3 — UnrollCycle", BLACK, HDR, True)],
          space_before=8, space_after=2)
    _para(tf, [("   Premise(s):  ", BLACK, BODY, False),
               ("((R), s, 1) ∈ C", blue, BODY, False),
               ("  and  ", BLACK, BODY, False),
               ("1 < len(s)", blue, BODY, False)])
    _para(tf, [("   ⟹  ", BLACK, BODY, False),
               ("(s[0], s[1]) ∈ R⁺", blue, BODY, False)])

    # Δ state boxes on the right (mirroring slide 3)
    _state_box(slide, "Δ₀ = ", "{len(s) > 1, minimal(s)}",
               Inches(7.6), Inches(0.95), Inches(4.9), Inches(0.6))
    _state_box(slide, "Δ₀′ = ", "{len(s) > 1, minimal(s)}",
               Inches(7.6), Inches(2.0), Inches(4.9), Inches(0.6))
    _state_box(slide, "Δ₁ = ", "{1 < len(s), (s[0], s[1]) ∈ R⁺}",
               Inches(7.6), Inches(3.15), Inches(5.2), Inches(0.6))

    # Step 4 — explore both cases (lower half, split by a divider)
    s4 = slide.shapes.add_textbox(Inches(0.6), Inches(4.15), Inches(12.2),
                                  Inches(0.9))
    stf = s4.text_frame
    stf.word_wrap = True
    _para(stf, [("Step 4 — SplitCycleLen", BLACK, 20, True)],
          first=True, space_after=2)
    _para(stf, [("   Premise(s):  ", BLACK, 16, False),
                ("((R), s, 2) ∈ C", blue, 16, False),
                ("       ⟹   explore both cases:", BLACK, 16, False)])

    div = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(6.67), Inches(5.2),
                                     Inches(6.67), Inches(7.25))
    div.line.color.rgb = GRAY
    div.line.width = Pt(1.0)

    # left: case 2 < len(s)
    lc = slide.shapes.add_textbox(Inches(0.6), Inches(5.25), Inches(5.9),
                                  Inches(2.1))
    ltf = lc.text_frame
    ltf.word_wrap = True
    _para(ltf, [("Case  ", BLACK, 17, True), ("2 < len(s)", BLUE, 17, True),
                (" :", BLACK, 17, True)], first=True, space_after=5)
    _para(ltf, [("Δ = Δ ∪ {2 < len(s)}", blue, 16, False)])
    _para(ltf, [("continue: Steps 5–7  →  (next slide)", GRAY, 15, False)],
          space_before=3)

    # right: case 2 = len(s) -> contradiction -> unsat
    rc = slide.shapes.add_textbox(Inches(6.95), Inches(5.25), Inches(6.05),
                                  Inches(2.1))
    rtf = rc.text_frame
    rtf.word_wrap = True
    _para(rtf, [("Case  ", BLACK, 17, True), ("2 = len(s)", BLUE, 17, True),
                (" :", BLACK, 17, True)], first=True, space_after=5)
    _para(rtf, [("Δ = Δ ∪ {2 = len(s), s[0] = s[1]}", blue, 16, False)])
    _para(rtf, [("       ", BLACK, 13, False),
                ("(s[0] = s[1] by definition of minimal)", GRAY, 13, False)])
    _para(rtf, [("⟹  ", BLACK, 16, False),
                ("(s[0], s[0]) ∈ R⁺", blue, 16, False),
                ("   contradicts  ∀x. ¬((x,x) ∈ R⁺)", BLACK, 14, False)])
    _para(rtf, [("⟹  UNSAT — close this case", RED, 17, True)],
          space_before=4)
    return slide


def add_ex2_case_slide():
    slide = prs.slides.add_slide(BLANK)
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1),
                                 Inches(0.7))
    tr = t.text_frame.paragraphs[0].add_run()
    tr.text = "Example 2: case  2 < len(s)"
    tr.font.size = Pt(30)
    tr.font.bold = True
    tr.font.color.rgb = BLACK

    note = slide.shapes.add_textbox(Inches(0.7), Inches(0.92), Inches(12.0),
                                    Inches(0.4))
    nr = note.text_frame.paragraphs[0].add_run()
    nr.text = "From Step 4 (case 2 < len(s)):  Δ ⊇ {2 < len(s), " \
              "(s[0], s[1]) ∈ R⁺}"
    nr.font.size = Pt(16)
    nr.font.italic = True
    nr.font.color.rgb = GRAY

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(12.0),
                                   Inches(4.6))
    tf = box.text_frame
    tf.word_wrap = True
    HDR, BODY = 24, 20
    blue = RULE_BLUE
    # Step 5 — UnrollCycle
    _para(tf, [("Step 5 — UnrollCycle", BLACK, HDR, True)], first=True)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("((R), s, 2) ∈ C", blue, BODY, False),
               ("   and   ", BLACK, BODY, False),
               ("2 < len(s)", blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("Δ = Δ ∪ {(s[1], s[2]) ∈ R⁺}", blue, BODY, False)])
    # Step 6 — JoinUp
    _para(tf, [("Step 6 — JoinUp", BLACK, HDR, True)], space_before=12)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("(s[0], s[1]) ∈ R⁺", blue, BODY, False),
               ("   and   ", BLACK, BODY, False),
               ("(s[1], s[2]) ∈ R⁺", blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("Δ = Δ ∪ {(s[0], s[2]) ∈ R⁺}", blue, BODY, False)])
    # Step 7 — ContrMinimal
    _para(tf, [("Step 7 — ContrMinimal", BLACK, HDR, True)], space_before=12)
    _para(tf, [("    Premise(s):   ", BLACK, BODY, False),
               ("((R), s, 3) ∈ C,   0 ≤ 0 < 2 < 3,   (s[0], s[2]) ∈ R⁺",
                blue, BODY, False)])
    _para(tf, [("    ⟹   ", BLACK, BODY, False),
               ("UNSAT — close this case", RED, BODY, True)])

    # both cases closed -> unsat
    mb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(2.2), Inches(6.3), Inches(8.9),
                                Inches(0.85))
    mb.fill.solid()
    mb.fill.fore_color.rgb = WHITE
    mb.line.color.rgb = RED
    mb.line.width = Pt(2.5)
    mb.shadow.inherit = False
    mtf = mb.text_frame
    mtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    mp = mtf.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "Both cases of Step 4 close  ⟹  Example 2 is unsat"
    mr.font.size = Pt(22)
    mr.font.bold = True
    mr.font.color.rgb = RED
    return slide


add_ex2_instance_slide()     # slide 6
add_ex2_derivation_slide()   # slide 7
add_ex2_case_slide()         # slide 8


# ===== Example 3: SC litmus test (§4.3.1) ==================================
ORANGE = RGBColor(0xE6, 0x51, 0x00)
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
PO, RF, FR, MO = BLUE, GREEN, ORANGE, PURPLE  # relation colors


def _slide_title(slide, text, size=30):
    t = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.2),
                                 Inches(0.7))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = BLACK
    return t


def add_ex3_instance_slide():
    slide = prs.slides.add_slide(BLANK)
    _slide_title(slide, "Example 3: SC litmus test  (§4.3.1)")
    note = slide.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(12.0),
                                    Inches(0.4))
    nr = note.text_frame.paragraphs[0].add_run()
    nr.text = ("Do the one-thread SC litmus-test templates forbid every "
               "MCM-violating execution?   →   expect unsat")
    nr.font.size = Pt(15)
    nr.font.italic = True
    nr.font.color.rgb = GRAY

    # legend of the four relations
    leg = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(12.0),
                                   Inches(0.4))
    lp = leg.text_frame.paragraphs[0]
    for txt, col in [("po", PO), (" program order      ", BLACK),
                     ("rf", RF), (" reads-from      ", BLACK),
                     ("fr", FR), (" from-reads      ", BLACK),
                     ("mo", MO), (" modification order", BLACK)]:
        r = lp.add_run()
        r.text = txt
        r.font.size = Pt(15)
        r.font.bold = (col is not BLACK)
        r.font.color.rgb = col

    box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12.0),
                                   Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    B = 19
    blue = RULE_BLUE
    _para(tf, [("Litmus-test templates (irreflexive):", BLACK, 18, True)],
          first=True, space_after=3)
    _para(tf, [("    irreflexive(po⁺),  irreflexive(rf),  irreflexive(fr),  "
                "irreflexive(mo)", blue, B, False)])
    _para(tf, [("    irreflexive(po⁺ ⋈ rf),  irreflexive(po⁺ ⋈ fr),  "
                "irreflexive(po⁺ ⋈ mo)", blue, B, False)])
    _para(tf, [("Single thread:", BLACK, 18, True)],
          space_before=10, space_after=3)
    _para(tf, [("    is_singleton(Thread)", blue, B, False)])
    _para(tf, [("    thread(x) = thread(y)  ⟺  x = y  ∨  (x,y) ∈ po⁺  ∨  "
                "(y,x) ∈ po⁺", blue, B, False)])
    _para(tf, [("Identity:", BLACK, 18, True), ("   ∀a. (a, a) ∈ iden",
                blue, B, False)], space_before=10)
    _para(tf, [("Goal:", BLACK, 18, True),
               ("   ¬ acyclic(po ∪ rf ∪ fr ∪ mo)", blue, B, False)],
          space_before=10)
    return slide


def add_ex3_setup_slide():
    slide = prs.slides.add_slide(BLANK)
    _slide_title(slide, "Example 3: derivation — setup & first unroll")
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12.2),
                                   Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True
    H, B = 21, 17
    blue = RULE_BLUE
    _para(tf, [("Step 1 — Saturate  (ElimTClos)", BLACK, H, True)],
          first=True, space_after=2)
    _para(tf, [("    ⟹  ", BLACK, B, False),
               ("rf⁺ ≈ rf,   fr⁺ ≈ fr", blue, B, False),
               ("    (rf, fr relate disjoint sets)", GRAY, B - 3, False)])
    _para(tf, [("Step 2 — skipped", GRAY, H - 2, True),
               ("  (no new constants yet)", GRAY, B - 3, False)],
          space_before=8, space_after=2)
    _para(tf, [("Step 3 — InstCycle", BLACK, H, True)],
          space_before=8, space_after=2)
    _para(tf, [("    Premise(s):  ", BLACK, B, False),
               ("¬acyclic(po ∪ rf ∪ fr ∪ mo)", blue, B, False)])
    _para(tf, [("    ⟹  ", BLACK, B, False),
               ("C ∪= ((po, rf, fr, mo), s, 1);   "
                "Δ ∪= {len(s) > 1, minimal((po, rf, fr, mo), s)}",
                blue, B, False)])
    _para(tf, [("Step 3 — SplitCycleLen  (cnt = 1)", BLACK, H, True)],
          space_before=8, space_after=2)
    _para(tf, [("    1 ≈ len(s) :  ", GRAY, B, True),
               ("unsat — conflicts len(s) > 1   (close)", GRAY, B, False)])
    _para(tf, [("    1 < len(s) :  ", BLACK, B, True),
               ("continue", blue, B, False)])
    _para(tf, [("Step 3 — UnrollCycle  (cnt = 1)", BLACK, H, True)],
          space_before=8, space_after=2)
    _para(tf, [("    ⟹  (s[1], s[2]) ∈  ", BLACK, B, False),
               ("po⁺", PO, B, True), ("  ‖  ", BLACK, B, True),
               ("rf⁺", RF, B, True), ("  ‖  ", BLACK, B, True),
               ("fr⁺", FR, B, True), ("  ‖  ", BLACK, B, True),
               ("mo⁺", MO, B, True),
               ("     (4 cases; explore ", GRAY, B - 3, False),
               ("po⁺", PO, B - 3, True),
               (", others similar)", GRAY, B - 3, False)])
    return slide


def add_ex3_po_slide():
    slide = prs.slides.add_slide(BLANK)
    _slide_title(slide, "Example 3: case  (s[1], s[2]) ∈ po⁺")
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12.2),
                                   Inches(0.8))
    stf = sub.text_frame
    stf.word_wrap = True
    _para(stf, [("Step — SplitCycleLen  (cnt = 2)", BLACK, 21, True)],
          first=True, space_after=2)
    _para(stf, [("    Premise(s):  ", BLACK, 16, False),
                ("((po, rf, fr, mo), s, 2) ∈ C", RULE_BLUE, 16, False),
                ("       ⟹   explore both cases:", BLACK, 16, False)])

    div = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(6.67), Inches(2.05),
                                     Inches(6.67), Inches(7.25))
    div.line.color.rgb = GRAY
    div.line.width = Pt(1.0)
    blue = RULE_BLUE

    # left: 2 ≈ len(s) -> contradiction
    lc = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(5.95),
                                  Inches(5.0))
    ltf = lc.text_frame
    ltf.word_wrap = True
    _para(ltf, [("Case  ", BLACK, 18, True), ("2 ≈ len(s)", BLUE, 18, True),
                (" :", BLACK, 18, True)], first=True, space_after=5)
    _para(ltf, [("s[1] = s[2]   ", blue, 16, False),
                ("(by minimal)", GRAY, 13, False)])
    _para(ltf, [("⟹  (s[1], s[1]) ∈ po⁺", blue, 16, False)])
    _para(ltf, [("step 2: instantiate iden at s[1]", BLACK, 15, False)],
          space_before=4)
    _para(ltf, [("⟹  (s[1], s[1]) ∈ iden", blue, 16, False)])
    _para(ltf, [("irreflexive(po⁺):  po⁺ ⊓ iden ≈ ∅", blue, 16, False)],
          space_before=4)
    _para(ltf, [("⟹  UNSAT — close this case", RED, 17, True)],
          space_before=4)

    # right: 2 < len(s) -> unroll again
    rc = slide.shapes.add_textbox(Inches(6.95), Inches(2.1), Inches(6.0),
                                  Inches(5.0))
    rtf = rc.text_frame
    rtf.word_wrap = True
    _para(rtf, [("Case  ", BLACK, 18, True), ("2 < len(s)", BLUE, 18, True),
                (" :", BLACK, 18, True)], first=True, space_after=5)
    _para(rtf, [("continue  ⟹  UnrollCycle  (cnt = 2)", BLACK, 16, False)])
    _para(rtf, [("(s[2], s[3]) ∈  ", BLACK, 16, False),
                ("po⁺", PO, 16, True), (" ‖ ", BLACK, 16, True),
                ("rf⁺", RF, 16, True), (" ‖ ", BLACK, 16, True),
                ("fr⁺", FR, 16, True), (" ‖ ", BLACK, 16, True),
                ("mo⁺", MO, 16, True)], space_before=3)
    _para(rtf, [("explore  ", GRAY, 15, False),
                ("(s[2], s[3]) ∈ rf⁺", RF, 15, True),
                ("  (others similar)", GRAY, 15, False)], space_before=3)
    _para(rtf, [("→  next slide", GRAY, 15, False)], space_before=3)
    return slide


def add_ex3_rf_slide():
    slide = prs.slides.add_slide(BLANK)
    _slide_title(slide, "Example 3: case  (s[2], s[3]) ∈ rf⁺")
    # current Δ
    draw_delta_box_multi(
        slide, "Δ",
        [("rf⁺ ≈ rf", GREEN), ("2 < len(s)", BLACK),
         ("(s[1], s[2]) ∈ po⁺", PO), ("(s[2], s[3]) ∈ rf⁺", RF)],
        BLACK, Inches(9.0), Inches(1.0), Inches(3.8), Inches(2.0),
        hsize=18, bsize=13)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(8.1),
                                   Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    H, B = 20, 16
    blue = RULE_BLUE
    _para(tf, [("Step 1 — Join Up", BLACK, H, True)], first=True, space_after=2)
    _para(tf, [("    ⟹  ", BLACK, B, False),
               ("(s[1], s[3]) ∈ (po⁺ ⋈ rf)", blue, B, False)])
    _para(tf, [("Step 2 — instantiate thread at s[1], s[3]  (single thread)",
                BLACK, H, True)], space_before=8, space_after=2)
    _para(tf, [("    thread(s[1]) = thread(s[3])  ⟹", blue, B, False)])
    _para(tf, [("      s[1] = s[3]  ∨  (s[1],s[3]) ∈ po⁺  ∨  "
                "(s[3],s[1]) ∈ po⁺", blue, B, False),
               ("   (split 3)", GRAY, B - 3, False)])

    # three thread cases
    cbox = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12.2),
                                    Inches(3.0))
    ctf = cbox.text_frame
    ctf.word_wrap = True
    CB = 16
    _para(ctf, [("①  s[1] = s[3] :  ", BLACK, CB, True),
                ("with (s[1],s[3]) ∈ (po⁺ ⋈ rf)  ⟹  contradicts "
                 "irreflexive(po⁺ ⋈ rf).   ", blue, CB, False),
                ("UNSAT", RED, CB, True)], first=True, space_after=8)
    _para(ctf, [("②  (s[1], s[3]) ∈ po⁺ :  ", BLACK, CB, True),
                ("ContrMinimal — non-consecutive elements related by a "
                 "transitive closure  ⟹  smaller cycle.   ", blue, CB, False),
                ("UNSAT", RED, CB, True)], space_after=8)
    _para(ctf, [("③  (s[3], s[1]) ∈ po⁺ :  ", BLACK, CB, True),
                ("TClos Up II ⟹ (s[3],s[2]) ∈ po⁺;   Join Up ⟹ "
                 "(s[3],s[3]) ∈ (po⁺ ⋈ rf)  ⟹  contradicts "
                 "irreflexive(po⁺ ⋈ rf).   ", blue, CB, False),
                ("UNSAT", RED, CB, True)])

    mb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(2.0), Inches(6.45), Inches(9.3),
                                Inches(0.8))
    mb.fill.solid()
    mb.fill.fore_color.rgb = WHITE
    mb.line.color.rgb = RED
    mb.line.width = Pt(2.5)
    mb.shadow.inherit = False
    mb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    mp = mb.text_frame.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "All other cases similar  ⟹  Example 3 is unsat"
    mr.font.size = Pt(20)
    mr.font.bold = True
    mr.font.color.rgb = RED
    return slide


add_ex3_instance_slide()   # slide 9
add_ex3_setup_slide()      # slide 10
add_ex3_po_slide()         # slide 11
add_ex3_rf_slide()         # slide 12

OUT = "cycle_slides.pptx"
prs.save(OUT)
print("wrote", OUT)
